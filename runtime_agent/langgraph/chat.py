import asyncio
import shutil
import sqlite3
import traceback
import boto3
import os
import json
import re
import uuid
import base64
import info 
import utils
import bedrock_data_retention
import langgraph_agent
import mcp_config
import skill

from langchain_core.documents import Document
from urllib import parse
from io import BytesIO
from PIL import Image
from langchain_aws import ChatBedrock
from langchain_aws import ChatBedrockConverse
from langchain_openai import ChatOpenAI
from botocore.config import Config
from botocore.exceptions import ClientError
from langchain_core.prompts import MessagesPlaceholder, ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.messages import HumanMessage, AIMessage, ToolMessage, AIMessageChunk
from langchain.mcp import MCPAdapter
from tavily_tool_interceptor import wrap_tavily_mcp_tools
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langgraph.checkpoint.memory import MemorySaver

try:
    from langgraph.checkpoint.sqlite.aio import AsyncSqliteSaver
    SQLITE_CHECKPOINTER_AVAILABLE = True
except ImportError:
    AsyncSqliteSaver = None  # type: ignore[misc, assignment]
    SQLITE_CHECKPOINTER_AVAILABLE = False

import logging
import sys

logging.basicConfig(
    level=logging.INFO,  # Default to INFO level
    format='%(filename)s:%(lineno)d | %(message)s',
    handlers=[
        logging.StreamHandler(sys.stderr)
    ]
)
logger = logging.getLogger("chat")

workingDir = os.path.dirname(os.path.abspath(__file__))
config_path = os.path.join(workingDir, "config.json")

reasoning_mode = 'Disable'
debug_messages = []  # List to store debug messages

config = utils.load_config()
print(f"config: {config}")

projectName = config.get("projectName", "es")
bedrock_region = config.get("region", "ap-northeast-2")

accountId = config.get("accountId")
knowledge_base_name = config.get("knowledge_base_name")
s3_bucket = config.get("s3_bucket")
s3_prefix = "docs"
s3_image_prefix = "images"

path = config.get('sharing_url', '')
doc_prefix = "docs/"

model_name = "Claude 5.0 Sonnet"
model_type = "claude"
models = info.get_model_info(model_name)
model_id = models[0]["model_id"]
model_type = models[0]["model_type"]
debug_mode = "Enable"
user_id = "agent"
guardrail_enabled = True
memory_enabled = False
memory_id = None
actor_id = None
session_id = None

def update(userId, modelName, debugMode, guardrailEnabled=None, memoryEnabled=None):    
    global model_name, model_id, model_type, debug_mode, reasoning_mode
    global models, user_id, guardrail_enabled, memory_enabled

    if userId != user_id:
        user_id = userId
        logger.info(f"user_id: {user_id}")
    # Isolate artifacts + user skills under {SESSION_STORAGE_DIR}/{user_id}/
    # and ensure {user_id}/skills.list exists (create if missing).
    langgraph_agent.set_user_workspace(user_id)
    skill.set_user_workspace(user_id)

    if model_name != modelName:
        model_name = modelName
        logger.info(f"model_name: {model_name}")
        
        models = info.get_model_info(model_name)
        model_id = models[0]["model_id"]
        model_type = models[0]["model_type"]
                                
    if debug_mode != debugMode:
        debug_mode = debugMode        
        logger.info(f"debug_mode: {debug_mode}")

    if guardrailEnabled is not None and guardrail_enabled != guardrailEnabled:
        guardrail_enabled = guardrailEnabled
        logger.info(f"guardrail_enabled: {guardrail_enabled}")

    if memoryEnabled is not None and memory_enabled != memoryEnabled:
        memory_enabled = memoryEnabled
        logger.info(f"memory_enabled: {memory_enabled}")

def _guardrail_config() -> dict | None:
    if not guardrail_enabled:
        return None
    runtime_config = config or utils.load_config()
    guardrail_id = runtime_config.get("guardrail_id")
    if not guardrail_id:
        return None
    return {
        "guardrailIdentifier": guardrail_id,
        "guardrailVersion": runtime_config.get("guardrail_version", "DRAFT"),
        "trace": "enabled",
    }


def uses_converse_guardrail() -> bool:
    return bool(_guardrail_config() and model_type in ("claude", "nova"))


def check_input_guardrail(text: str) -> tuple[bool, str]:
    """Return (blocked, message). When blocked, message is the guardrail response."""
    guardrail_cfg = _guardrail_config()
    if not guardrail_cfg or not text:
        return False, text

    try:
        client = boto3.client("bedrock-runtime", region_name=bedrock_region)
        response = client.apply_guardrail(
            guardrailIdentifier=guardrail_cfg["guardrailIdentifier"],
            guardrailVersion=guardrail_cfg["guardrailVersion"],
            source="INPUT",
            content=[{"text": {"text": text}}],
        )
        if response.get("action") == "GUARDRAIL_INTERVENED":
            logger.info("Guardrail blocked user input")
            for output in response.get("outputs", []):
                text_output = output.get("text", {})
                if text_output.get("text"):
                    return True, text_output["text"]
            return (
                True,
                "요청이 안전 정책에 의해 차단되었습니다. "
                "성적 표현 또는 프롬프트 공격이 감지되었습니다.",
            )
    except ClientError as e:
        logger.error(f"apply_guardrail failed: {e}")
    except Exception as e:
        logger.error(f"apply_guardrail failed: {e}")
    return False, text

SESSION_STORAGE_DIR = os.environ.get(
    "SESSION_STORAGE_DIR",
    "/mnt/workspace" if os.path.isdir("/mnt/workspace") else os.path.join(workingDir, ".session_storage"),
)
LEGACY_CHECKPOINT_DB = os.path.join(SESSION_STORAGE_DIR, "langgraph_checkpoints.sqlite")

checkpointer = MemorySaver()
_memory_checkpointer: MemorySaver | None = None
_sqlite_checkpointer = None
_sqlite_checkpointer_cm = None
_active_checkpoint_session = None
_current_checkpoint_session_id: str | None = None
_checkpointer_init_lock = asyncio.Lock()
SQLITE_BUSY_TIMEOUT_MS = 5000
_SETUP_MAX_ATTEMPTS = 5
_SETUP_RETRY_BASE_SEC = 0.25


def _runtime_session_id() -> str | None:
    try:
        from bedrock_agentcore.runtime.context import BedrockAgentCoreContext

        return BedrockAgentCoreContext.get_session_id()
    except Exception:
        return None


def set_checkpoint_session_id(session_id: str | None) -> None:
    """Bind the current request to a task's runtime_session_id."""
    global _current_checkpoint_session_id
    _current_checkpoint_session_id = session_id.strip() if session_id else None


def _checkpoint_session_id() -> str | None:
    """Resolve the active task session for checkpoint DB and thread_id."""
    if _current_checkpoint_session_id:
        return _current_checkpoint_session_id
    return _runtime_session_id()


def get_checkpoint_db_path() -> str:
    """Working SQLite path on local disk (avoids session-storage locking during runtime)."""
    session_id = _checkpoint_session_id()
    if session_id:
        local_dir = os.path.join("/tmp", "langgraph-checkpoints", session_id)
        os.makedirs(local_dir, exist_ok=True)
        return os.path.join(local_dir, "langgraph_checkpoints.sqlite")
    return LEGACY_CHECKPOINT_DB


def get_persistent_checkpoint_db_path() -> str:
    """Durable checkpoint path on AgentCore session storage (/mnt/workspace), per task."""
    session_id = _checkpoint_session_id()
    if session_id:
        checkpoint_dir = os.path.join(SESSION_STORAGE_DIR, "checkpoints", session_id)
        os.makedirs(checkpoint_dir, exist_ok=True)
        return os.path.join(checkpoint_dir, "langgraph_checkpoints.sqlite")
    return LEGACY_CHECKPOINT_DB


def _restore_from_session_storage(working_db: str) -> None:
    """Copy durable checkpoint from session storage into the local working DB."""
    persistent = get_persistent_checkpoint_db_path()
    if working_db == persistent:
        return

    if not _checkpoint_db_ready(persistent):
        if os.path.isfile(persistent):
            logger.warning(
                f"Persistent checkpoint empty, skip restore: {persistent} "
                f"(size={os.path.getsize(persistent)})"
            )
        elif os.path.isdir(SESSION_STORAGE_DIR):
            try:
                entries = os.listdir(SESSION_STORAGE_DIR)
            except OSError as exc:
                entries = [f"<listdir failed: {exc}>"]
            logger.warning(
                f"No persistent checkpoint at {persistent}; "
                f"session storage dir {SESSION_STORAGE_DIR} contents={entries}"
            )
        else:
            logger.warning(
                f"Session storage unavailable, skip restore: {SESSION_STORAGE_DIR} "
                f"(expected {persistent})"
            )
        return

    if _checkpoint_db_ready(working_db):
        if os.path.getmtime(persistent) <= os.path.getmtime(working_db):
            logger.info(
                f"Working checkpoint is newer, skip restore: working={working_db}, "
                f"persistent={persistent}"
            )
            return

    os.makedirs(os.path.dirname(working_db), exist_ok=True)
    shutil.copy2(persistent, working_db)
    for suffix in ("-wal", "-shm"):
        src = persistent + suffix
        if os.path.isfile(src):
            shutil.copy2(src, working_db + suffix)
    logger.info(f"Restored checkpoint DB from session storage: {persistent} -> {working_db}")


async def persist_checkpoint_to_session_storage() -> None:
    """Flush working checkpoint to session storage so history survives microVM stop/resume."""
    if _sqlite_checkpointer is None:
        return

    working_db = get_checkpoint_db_path()
    persistent = get_persistent_checkpoint_db_path()
    if working_db == persistent or not _checkpoint_db_ready(working_db):
        return

    os.makedirs(SESSION_STORAGE_DIR, exist_ok=True)

    try:
        await _sqlite_checkpointer.conn.execute("PRAGMA wal_checkpoint(TRUNCATE)")
        await _sqlite_checkpointer.conn.commit()

        def _copy_checkpoint_files():
            shutil.copy2(working_db, persistent)
            for suffix in ("-wal", "-shm"):
                src = working_db + suffix
                dst = persistent + suffix
                if os.path.isfile(src):
                    shutil.copy2(src, dst)
                elif os.path.isfile(dst):
                    os.remove(dst)

        await asyncio.to_thread(_copy_checkpoint_files)
        logger.info(f"Checkpoint persisted to session storage: {persistent}")
    except Exception as exc:
        logger.warning(f"Failed to persist checkpoint to session storage: {exc}")


def _is_sqlite_locked_error(exc: BaseException) -> bool:
    if isinstance(exc, sqlite3.OperationalError):
        return "locked" in str(exc).lower()
    cause = getattr(exc, "__cause__", None)
    return isinstance(cause, sqlite3.OperationalError) and "locked" in str(cause).lower()


def _checkpoint_db_ready(checkpoint_db: str) -> bool:
    return os.path.isfile(checkpoint_db) and os.path.getsize(checkpoint_db) > 0


async def _configure_sqlite_connection(conn) -> None:
    await conn.execute(f"PRAGMA busy_timeout={SQLITE_BUSY_TIMEOUT_MS}")
    await conn.commit()


async def _open_existing_sqlite_checkpointer(checkpoint_db: str):
    import aiosqlite

    for attempt in range(1, _SETUP_MAX_ATTEMPTS + 1):
        conn = None
        try:
            conn = await aiosqlite.connect(
                checkpoint_db,
                timeout=SQLITE_BUSY_TIMEOUT_MS / 1000,
            )
            await _configure_sqlite_connection(conn)
            saver = AsyncSqliteSaver(conn)
            saver.is_setup = True
            return saver
        except Exception as exc:
            if conn is not None:
                await conn.close()
            if not _is_sqlite_locked_error(exc):
                raise
            if attempt == _SETUP_MAX_ATTEMPTS:
                raise
            delay = _SETUP_RETRY_BASE_SEC * (2 ** (attempt - 1))
            logger.warning(
                f"SQLite open locked (attempt {attempt}/{_SETUP_MAX_ATTEMPTS}), "
                f"retrying in {delay:.2f}s"
            )
            await asyncio.sleep(delay)


async def _create_sqlite_checkpointer(checkpoint_db: str):
    import aiosqlite

    for attempt in range(1, _SETUP_MAX_ATTEMPTS + 1):
        conn = None
        try:
            conn = await aiosqlite.connect(
                checkpoint_db,
                timeout=SQLITE_BUSY_TIMEOUT_MS / 1000,
            )
            await _configure_sqlite_connection(conn)
            saver = AsyncSqliteSaver(conn)
            await saver.setup()
            return saver
        except Exception as exc:
            if conn is not None:
                await conn.close()
            if not _is_sqlite_locked_error(exc):
                raise
            if attempt == _SETUP_MAX_ATTEMPTS:
                raise
            delay = _SETUP_RETRY_BASE_SEC * (2 ** (attempt - 1))
            logger.warning(
                f"SQLite checkpointer locked (attempt {attempt}/{_SETUP_MAX_ATTEMPTS}), "
                f"retrying in {delay:.2f}s"
            )
            await asyncio.sleep(delay)


async def _close_sqlite_checkpointer() -> None:
    global _sqlite_checkpointer, _sqlite_checkpointer_cm

    if _sqlite_checkpointer is None:
        return

    conn = getattr(_sqlite_checkpointer, "conn", None)
    if conn is not None:
        try:
            await conn.close()
        except Exception as exc:
            logger.warning(f"Failed to close sqlite checkpointer connection: {exc}")

    _sqlite_checkpointer = None
    _sqlite_checkpointer_cm = None


def _memory_checkpointer() -> MemorySaver:
    """Reuse a single in-memory saver when SQLite is unavailable."""
    global _memory_checkpointer, checkpointer
    if _memory_checkpointer is None:
        _memory_checkpointer = MemorySaver()
    checkpointer = _memory_checkpointer
    return _memory_checkpointer


async def ensure_checkpointer():
    """Initialize AsyncSqliteSaver on local disk (per task runtime_session_id)."""
    global checkpointer, _sqlite_checkpointer, _sqlite_checkpointer_cm, _active_checkpoint_session

    session_id = _checkpoint_session_id()
    checkpoint_db = get_checkpoint_db_path()

    if _sqlite_checkpointer is not None and _active_checkpoint_session == session_id:
        checkpointer = _sqlite_checkpointer
        return checkpointer

    if not SQLITE_CHECKPOINTER_AVAILABLE:
        logger.info("Using in-memory checkpointer (langgraph-checkpoint-sqlite not installed)")
        return _memory_checkpointer()

    async with _checkpointer_init_lock:
        if _sqlite_checkpointer is not None and _active_checkpoint_session == session_id:
            checkpointer = _sqlite_checkpointer
            return checkpointer

        if _sqlite_checkpointer is not None and _active_checkpoint_session != session_id:
            logger.info(
                "Switching checkpointer session: %s -> %s",
                _active_checkpoint_session,
                session_id,
            )
            await _close_sqlite_checkpointer()

        _active_checkpoint_session = session_id

        try:
            _restore_from_session_storage(checkpoint_db)
            if _checkpoint_db_ready(checkpoint_db):
                _sqlite_checkpointer = await _open_existing_sqlite_checkpointer(checkpoint_db)
                logger.info(
                    "SQLite checkpointer opened (existing): session=%s path=%s",
                    session_id,
                    checkpoint_db,
                )
            else:
                _sqlite_checkpointer = await _create_sqlite_checkpointer(checkpoint_db)
                logger.info(
                    "SQLite checkpointer initialized: session=%s path=%s",
                    session_id,
                    checkpoint_db,
                )
        except Exception as exc:
            logger.error(
                f"SQLite checkpointer unavailable ({exc}); falling back to MemorySaver"
            )
            return _memory_checkpointer()

        checkpointer = _sqlite_checkpointer
        return checkpointer


def _thread_scope(runtime_session_id: str | None) -> str:
    """Isolate LangGraph checkpoint threads per application task."""
    if runtime_session_id:
        return runtime_session_id

    import hashlib

    payload = json.dumps({"user_id": user_id}, sort_keys=True)
    digest = hashlib.sha256(payload.encode()).hexdigest()[:12]
    logger.warning(
        "runtime_session_id missing; falling back to shared thread scope for user %s",
        user_id,
    )
    return f"{user_id}:{digest}:default"


selected_chat = 0


def is_fable_model(model_id: str | None = None) -> bool:
    if not model_id:
        if not models:
            return False
        model_id = models[0].get("model_id", "")
    return "fable" in model_id.lower()


def uses_adaptive_thinking(model_id: str | None = None) -> bool:
    if not model_id:
        if not models:
            return False
        model_id = models[0].get("model_id", "")
    mid = model_id.lower()
    return (
        "fable" in mid
        or "claude-sonnet-5" in mid
        or "claude-5-sonnet" in mid
        or "claude-opus-5" in mid
        or "claude-5-opus" in mid
    )


# Content block types that Bedrock Claude/Nova reject on replay (e.g. after
# OpenAI Responses API turns leave type='reasoning', or signed thinking is gone).
_BEDROCK_NON_REPLAYABLE_CONTENT_TYPES = frozenset({"thinking", "reasoning"})


def sanitize_adaptive_thinking_messages(messages: list) -> list:
    """Remove content blocks that cannot be replayed to Bedrock on later turns.

    Strips Anthropic ``thinking`` and OpenAI Responses ``reasoning`` blocks so
    model switches (e.g. GPT → Claude) do not hit ValidationException.
    """
    sanitized = []
    for msg in messages:
        if not isinstance(msg, AIMessage):
            sanitized.append(msg)
            continue

        content = msg.content
        if not isinstance(content, list):
            sanitized.append(msg)
            continue

        cleaned = [
            block for block in content
            if not (
                isinstance(block, dict)
                and block.get("type") in _BEDROCK_NON_REPLAYABLE_CONTENT_TYPES
            )
        ]
        if not cleaned:
            cleaned = ""
        elif (
            len(cleaned) == 1
            and isinstance(cleaned[0], dict)
            and cleaned[0].get("type") == "text"
        ):
            cleaned = cleaned[0].get("text", "")

        sanitized.append(
            AIMessage(
                content=cleaned,
                tool_calls=getattr(msg, "tool_calls", None) or [],
                additional_kwargs=getattr(msg, "additional_kwargs", {}),
                response_metadata=getattr(msg, "response_metadata", {}),
                usage_metadata=getattr(msg, "usage_metadata", None),
                id=getattr(msg, "id", None),
            )
        )
    return sanitized


def get_max_output_tokens(model_id: str = "") -> int:
    """Return the max output tokens based on the model ID."""
    mid = (model_id or "").lower()
    if is_fable_model(model_id):
        return 128000
    if "claude-sonnet-5" in mid or "claude-5-sonnet" in mid:
        return 128000
    if "claude-opus-5" in mid or "claude-5-opus" in mid:
        return 128000
    if "claude-opus-4-6" in mid:
        return 128000
    if "claude-opus-4-5" in mid:
        return 64000
    if "claude-opus-4" in mid or "claude-4-opus" in mid:
        return 32000
    if "claude-sonnet-4" in mid or "claude-4-sonnet" in mid or "claude-haiku-4" in mid:
        return 64000
    return 8192
    
def _build_openai_chat(profile: dict, max_output_tokens: int):
    """Build OpenAI-on-Bedrock chat model (Mantle Responses API or invoke_model)."""
    bedrock_region = profile["bedrock_region"]
    model_id = profile["model_id"]
    mantle_api = profile.get("mantle_api", "chat")

    if mantle_api == "responses":
        def bearer_token_provider() -> str:
            return bedrock_data_retention.get_bedrock_bearer_token(bedrock_region)

        return ChatOpenAI(
            model=model_id,
            api_key=bearer_token_provider,
            base_url=f"https://bedrock-mantle.{bedrock_region}.api.aws/openai/v1",
            use_responses_api=True,
            max_tokens=max_output_tokens,
        )

    boto3_bedrock = boto3.client(
        service_name="bedrock-runtime",
        region_name=bedrock_region,
        config=Config(
            retries={"max_attempts": 30},
            read_timeout=300,
        ),
    )
    chat = ChatBedrock(
        model_id=model_id,
        client=boto3_bedrock,
        model_kwargs={
            "max_tokens": max_output_tokens,
        },
        region_name=bedrock_region,
    )
    chat.streaming = False
    return chat

def get_chat():
    global model_type

    logger.info(f"models: {models}")

    profile = models[0]
    modelId = profile['model_id']
    model_type = profile['model_type']
    bedrock_region = profile['bedrock_region']
    if model_type == 'claude':
        maxOutputTokens = get_max_output_tokens(modelId)
    else:
        maxOutputTokens = 5120 # 5k

    logger.info(f"modelId: {modelId}, model_type: {model_type}, bedrock_region: {bedrock_region}")

    if is_fable_model(modelId):
        bedrock_data_retention.ensure_fable_data_retention(
            modelId,
            bedrock_region=bedrock_region,
        )

    if profile["model_type"] == "openai":
        return _build_openai_chat(profile, maxOutputTokens)

    guardrail_cfg = _guardrail_config()
    if guardrail_cfg and profile["model_type"] in ("claude", "nova"):
        boto3_bedrock = boto3.client(
            service_name="bedrock-runtime",
            region_name=bedrock_region,
            config=Config(
                retries={"max_attempts": 30},
                read_timeout=300,
            ),
        )
        converse_kwargs = {
            "model_id": modelId,
            "client": boto3_bedrock,
            "max_tokens": maxOutputTokens,
            "temperature": 0.1,
            "region_name": bedrock_region,
            "guardrail_config": guardrail_cfg,
        }
        if model_type == "claude":
            converse_kwargs["provider"] = "anthropic"
        converse_chat = ChatBedrockConverse(**converse_kwargs)
        converse_chat.streaming = False
        return converse_chat

    if profile['model_type'] == 'nova':
        STOP_SEQUENCE = '"\n\n<thinking>", "\n<thinking>", " <thinking>"'
    elif profile['model_type'] == 'claude':
        STOP_SEQUENCE = "\n\nHuman:"
    else:
        STOP_SEQUENCE = ""
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            },
            read_timeout=300
        )
    )

    parameters = {
        "max_tokens":maxOutputTokens,     
        "stop_sequences": [STOP_SEQUENCE]
    }

    chat_kwargs = {
        "model_id": modelId,
        "client": boto3_bedrock,
        "model_kwargs": parameters,
        "region_name": bedrock_region,
    }
    if model_type == "claude":
        chat_kwargs["provider"] = "anthropic"

    chat = ChatBedrock(**chat_kwargs)
    
    return chat

reference_docs = []

def upload_to_s3(file_bytes, file_name):
    """
    Upload a file to S3 under {prefix}/{user_id}/ and return the URL.
    """

    try:
        s3_client = boto3.client(
            service_name='s3',
            region_name=bedrock_region,
        )

        content_type = utils.get_contents_type(file_name)       
        logger.info(f"content_type: {content_type}") 

        if content_type.startswith("image/"):
            prefix = s3_image_prefix
        else:
            prefix = s3_prefix

        user_segment = utils.sanitize_user_path_segment(user_id)
        if user_segment:
            s3_key = f"{prefix}/{user_segment}/{file_name}"
            relative_url_path = f"{prefix}/{parse.quote(user_segment)}/{parse.quote(file_name)}"
        else:
            s3_key = f"{prefix}/{file_name}"
            relative_url_path = f"{prefix}/{parse.quote(file_name)}"
        
        user_meta = {  # user-defined metadata
            "content_type": content_type,
            "model_name": model_name
        }
        
        response = s3_client.put_object(
            Bucket=s3_bucket, 
            Key=s3_key, 
            ContentType=content_type,
            Metadata = user_meta,
            Body=file_bytes            
        )
        logger.info(f"upload response: {response}")

        url = path.rstrip("/") + "/" + relative_url_path if path else relative_url_path
        return url
    
    except Exception as e:
        err_msg = f"Error uploading to S3: {str(e)}"
        logger.info(f"{err_msg}")
        return None

def isKorean(text):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(text))
    # print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        # logger.info(f"Korean: {word_kor}")
        return True
    else:
        # logger.info(f"Not Korean:: {word_kor}")
        return False
    
def traslation(chat, text, input_language, output_language):
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags." 
        "Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        # print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        logger.info(f"error message: {err_msg}")     
        raise Exception ("Not able to request to LLM")

    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def get_summary(docs):    
    llm = get_chat()

    text = ""
    for doc in docs:
        text = text + doc
    
    if isKorean(text)==True:
        system = (
            "다음의 <article> tag안의 문장을 요약해서 500자 이내로 설명하세오."
        )
    else: 
        system = (
            "Here is pieces of article, contained in <article> tags. Write a concise summary within 500 characters."
        )
    
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | llm    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )
        
        summary = result.content
        logger.info(f"esult of summarization: {summary}")
    except Exception:
        err_msg = traceback.format_exc()
        logger.info(f"error message: {err_msg}") 
        raise Exception ("Not able to request to LLM")
    
    return summary


def _content_to_text(content: object) -> str:
    """Normalize LangChain/Bedrock message content to a plain string.

    Newer Claude/Bedrock multimodal responses often return ``content`` as a list
    of blocks (e.g. ``[{"type": "text", "text": "..."}]``).
    """
    if content is None:
        return ""
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts: list[str] = []
        for item in content:
            if isinstance(item, str):
                parts.append(item)
            elif isinstance(item, dict):
                text = item.get("text")
                if text:
                    parts.append(str(text))
            else:
                text = getattr(item, "text", None)
                if text:
                    parts.append(str(text))
        return "\n".join(parts).strip()
    return str(content).strip()


def _parse_result_tag(text: str) -> str:
    """Extract inner text from <result>...</result> when present."""
    if not text:
        return text
    start = text.find("<result>")
    end = text.find("</result>")
    if start != -1 and end != -1 and end > start:
        return text[start + 8 : end]
    return text


def summary_image(img_base64, instruction):      
    llm = get_chat()

    if instruction:
        logger.info(f"instruction: {instruction}")
        query = f"{instruction}. <result> tag를 붙여주세요. 한국어로 답변하세요."
        
    else:
        query = "이미지가 의미하는 내용을 풀어서 자세히 알려주세요. markdown 포맷으로 답변을 작성합니다."
    
    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    extracted_text = ""
    for attempt in range(5):
        logger.info(f"attempt: {attempt}")
        try: 
            result = llm.invoke(messages)
            
            extracted_text = _content_to_text(result.content)
            # print('summary from an image: ', extracted_text)
            break
        except Exception:
            err_msg = traceback.format_exc()
            logger.info(f"error message: {err_msg}")                    
            raise Exception ("Not able to request to LLM")
        
    return extracted_text

def extract_text(img_base64):    
    multimodal = get_chat()
    query = "텍스트를 추출해서 markdown 포맷으로 변환하세요. 원문의 언어를 그대로 유지하고 번역하지 마세요. <result> tag를 붙여주세요."
    
    extracted_text = ""
    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    for attempt in range(5):
        logger.info(f"attempt: {attempt}")
        try: 
            result = multimodal.invoke(messages)
            
            extracted_text = _content_to_text(result.content)
            # print('result of text extraction from an image: ', extracted_text)
            break
        except Exception:
            err_msg = traceback.format_exc()
            logger.info(f"error message: {err_msg}")                    
            # raise Exception ("Not able to request to LLM")
    
    logger.info(f"Extracted_text: {extracted_text}")
    if len(extracted_text)<10:
        extracted_text = "텍스트를 추출하지 못하였습니다."    

    return extracted_text

fileId = uuid.uuid4().hex
# print('fileId: ', fileId)

####################### LangChain #######################
# Image Summarization
#########################################################
def summarize_image(image_content: bytes, prompt: str) -> str:
    img = Image.open(BytesIO(image_content))
    
    width, height = img.size 
    logger.info(f"width: {width}, height: {height}, size: {width*height}")
    
    # 이미지 리사이징 및 크기 확인
    isResized = False
    max_size = 5 * 1024 * 1024  # 5MB in bytes
    
    # Initial resizing (based on pixel count)
    while(width*height > 2000000):  # Limit to approximately 2M pixels
        width = int(width/2)
        height = int(height/2)
        isResized = True
        logger.info(f"width: {width}, height: {height}, size: {width*height}")
    
    if isResized:
        img = img.resize((width, height))
    
    # Base64 size verification and additional resizing
    max_attempts = 5
    for attempt in range(max_attempts):
        buffer = BytesIO()
        img.save(buffer, format="PNG", optimize=True)
        img_bytes = buffer.getvalue()
        img_base64 = base64.b64encode(img_bytes).decode("utf-8")
        
        # Base64 size verification (actual transmission size)
        base64_size = len(img_base64.encode('utf-8'))
        logger.info(f"attempt {attempt + 1}: base64_size = {base64_size} bytes")
        
        if base64_size <= max_size:
            break
        else:
            # Resize smaller if still too large
            width = int(width * 0.8)
            height = int(height * 0.8)
            img = img.resize((width, height))
            logger.info(f"resizing to {width}x{height} due to size limit")
    
    if base64_size > max_size:
        logger.warning(f"Image still too large after {max_attempts} attempts: {base64_size} bytes")
        raise Exception(f"이미지 크기가 너무 큽니다. 5MB 이하의 이미지를 사용해주세요.")

    logger.info("이미지의 내용을 분석합니다.")
    result = summary_image(img_base64, prompt)
    if not isinstance(result, str):
        result = _content_to_text(result)

    summary = _parse_result_tag(result)
    logger.info(f"image summary: {summary}")
            
    return summary


def _file_name_from_ref(file_ref: str) -> str:
    """Extract a basename from a sharing URL or plain file name."""
    raw = (file_ref or "").strip()
    if not raw:
        return ""
    name = raw.rsplit("/", 1)[-1]
    return parse.unquote(name)


def _s3_key_from_file_ref(file_ref: str, *, default_prefix: str = s3_image_prefix) -> str:
    """Derive S3 object key from a sharing URL or bare file name.

    Prefer the ``images/{user_id}/...`` (or docs/) path embedded in the URL;
    otherwise fall back to ``{prefix}/{user_id}/{file_name}``.
    """
    raw = (file_ref or "").strip()
    if not raw:
        return ""

    raw = raw.split("?", 1)[0].split("#", 1)[0]

    for prefix in (s3_image_prefix, s3_prefix, "images", "docs"):
        marker = f"/{prefix}/"
        idx = raw.find(marker)
        if idx >= 0:
            return parse.unquote(raw[idx + 1 :])
        if raw.startswith(f"{prefix}/"):
            return parse.unquote(raw)

    file_name = _file_name_from_ref(file_ref)
    if not file_name:
        return ""

    user_segment = utils.sanitize_user_path_segment(user_id)
    if user_segment:
        return f"{default_prefix}/{user_segment}/{file_name}"
    return f"{default_prefix}/{file_name}"


def _workspace_ref_to_s3_key(file_ref: str) -> str | None:
    """Map /mnt/workspace/{user}/upload/x → agentcore-sessions/{user}/upload/x."""
    path = (file_ref or "").strip()
    marker = "/mnt/workspace/"
    if not path.startswith(marker):
        return None
    rel = path[len(marker) :].lstrip("/")
    if not rel or ".." in rel.split("/"):
        return None
    return f"agentcore-sessions/{rel}"


def _wait_for_workspace_mount_file(
    path: str,
    *,
    timeout_sec: float = 90.0,
    interval_sec: float = 1.0,
) -> bool:
    """Poll until ``path`` appears under /mnt/workspace (S3 Files lag)."""
    import time

    if not path.startswith("/mnt/workspace/"):
        return False
    deadline = time.monotonic() + max(0.0, timeout_sec)
    while True:
        if os.path.isfile(path) and os.path.getsize(path) > 0:
            logger.info(f"workspace mount file ready: {path}")
            return True
        if time.monotonic() >= deadline:
            logger.warning(
                f"workspace mount file not visible after {timeout_sec:.0f}s: {path}"
            )
            return False
        time.sleep(max(0.1, interval_sec))


def _load_workspace_file_bytes(file_ref: str) -> tuple[bytes | None, str]:
    """Load Load-files bytes from /mnt/workspace, waiting for mount if needed.

    Falls back to S3 API only if the mount never shows the file.
    """
    path = (file_ref or "").strip()
    if path.startswith("/mnt/workspace/") and os.path.isfile(path):
        with open(path, "rb") as f:
            data = f.read()
        logger.info(f"loaded workspace file from mount ({len(data)} bytes): {path}")
        return data, path

    if path.startswith("/mnt/workspace/"):
        logger.info(f"waiting for workspace mount file: {path}")
        if _wait_for_workspace_mount_file(path, timeout_sec=90.0, interval_sec=1.0):
            with open(path, "rb") as f:
                data = f.read()
            logger.info(
                f"loaded workspace file after wait ({len(data)} bytes): {path}"
            )
            return data, path

    s3_key = _workspace_ref_to_s3_key(path)
    if not s3_key or not s3_bucket:
        logger.warning(
            "workspace file unavailable on mount and no S3 key/bucket: path=%s key=%s",
            path,
            s3_key,
        )
        return None, path

    try:
        s3_client = boto3.client(service_name="s3", region_name=bedrock_region)
        logger.info(f"loading workspace file from s3://{s3_bucket}/{s3_key}")
        obj = s3_client.get_object(Bucket=s3_bucket, Key=s3_key)
        data = obj["Body"].read()
        logger.info(f"loaded workspace file from S3 ({len(data)} bytes): {s3_key}")
        try:
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "wb") as f:
                f.write(data)
        except Exception as cache_err:
            logger.warning(f"could not cache S3 file onto mount {path}: {cache_err}")
        return data, f"s3://{s3_bucket}/{s3_key}"
    except Exception:
        logger.error(
            "Failed to load workspace file from S3 key=%s: %s",
            s3_key,
            traceback.format_exc(),
        )
        return None, path


def _extract_text_from_docx_bytes(data: bytes) -> str:
    try:
        import docx

        document = docx.Document(BytesIO(data))
        parts: list[str] = []
        for p in document.paragraphs:
            if p.text and p.text.strip():
                parts.append(p.text)
        for table in document.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        text = "\n".join(parts).strip()
        if text:
            return text
    except Exception as e:
        logger.warning(f"python-docx extract failed, trying zip/xml: {e}")

    import zipfile
    from xml.etree import ElementTree as ET

    with zipfile.ZipFile(BytesIO(data)) as zf:
        xml_bytes = zf.read("word/document.xml")
    root = ET.fromstring(xml_bytes)
    paragraphs: list[str] = []
    for node in root.iter():
        if node.tag.endswith("}p") or node.tag == "p":
            runs = [
                (t.text or "")
                for t in node.iter()
                if (t.tag.endswith("}t") or t.tag == "t") and t.text
            ]
            line = "".join(runs).strip()
            if line:
                paragraphs.append(line)
    if paragraphs:
        return "\n".join(paragraphs)
    texts = [
        (node.text or "")
        for node in root.iter()
        if node.tag.endswith("}t") or node.tag == "t"
    ]
    return "\n".join(t for t in texts if t).strip()


def _extract_text_from_legacy_doc_bytes(data: bytes) -> str:
    """Best-effort extraction for Word 97-2003 (.doc)."""
    import subprocess
    import tempfile

    with tempfile.TemporaryDirectory() as tmp:
        doc_path = os.path.join(tmp, "input.doc")
        with open(doc_path, "wb") as f:
            f.write(data)

        for cmd in (
            ["antiword", doc_path],
            ["catdoc", "-w", doc_path],
        ):
            try:
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=60,
                    check=False,
                )
                text = (result.stdout or "").strip()
                if result.returncode == 0 and text:
                    return text
            except FileNotFoundError:
                continue
            except Exception as e:
                logger.warning("legacy .doc extractor %s failed: %s", cmd[0], e)

        # LibreOffice → txt
        try:
            result = subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "txt:Text",
                    "--outdir",
                    tmp,
                    doc_path,
                ],
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
            txt_path = os.path.join(tmp, "input.txt")
            if os.path.isfile(txt_path):
                with open(txt_path, "r", encoding="utf-8", errors="replace") as f:
                    text = f.read().strip()
                if text:
                    return text
            if result.returncode != 0:
                logger.warning(
                    "soffice .doc convert failed: %s",
                    (result.stderr or result.stdout or "")[:500],
                )
        except FileNotFoundError:
            pass
        except Exception as e:
            logger.warning("soffice .doc convert failed: %s", e)

    # Last resort: pull printable / UTF-16LE runs from the OLE binary.
    try:
        import re

        texts: list[str] = []
        for offset in (0, 1):
            for match in re.finditer(rb"(?:[\x20-\x7e]\x00){6,}", data[offset:]):
                try:
                    chunk = match.group().decode("utf-16le", errors="ignore").strip()
                except Exception:
                    continue
                if len(chunk) >= 6 and chunk not in texts:
                    texts.append(chunk)
        ascii_runs = re.findall(rb"[\x20-\x7e]{8,}", data)
        for raw in ascii_runs:
            chunk = raw.decode("ascii", errors="ignore").strip()
            if len(chunk) >= 8 and chunk not in texts:
                texts.append(chunk)
        joined = "\n".join(texts).strip()
        if len(joined) >= 40:
            return joined
    except Exception as e:
        logger.warning("legacy .doc binary string extract failed: %s", e)

    return ""


def _extract_text_from_pdf_bytes(data: bytes) -> str:
    try:
        import pdfplumber

        parts: list[str] = []
        with pdfplumber.open(BytesIO(data)) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text)
        text = "\n\n".join(parts).strip()
        if text:
            return text
    except Exception as e:
        logger.warning("pdfplumber bytes extract failed (%s); trying pypdf", e)

    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(data))
        parts = [(page.extract_text() or "") for page in reader.pages]
        return "\n\n".join(parts).strip()
    except Exception as e:
        raise RuntimeError(f"PDF 텍스트 추출 실패: {e}") from e


def _extract_text_from_document_bytes(data: bytes, file_type: str) -> str:
    """Best-effort text extraction for Load-files (bytes from mount or S3)."""
    if file_type in (
        "txt",
        "md",
        "markdown",
        "csv",
        "json",
        "py",
        "js",
        "ts",
        "tsx",
        "jsx",
        "html",
        "htm",
        "yml",
        "yaml",
        "xml",
        "rst",
    ):
        return data.decode("utf-8", errors="replace")

    if file_type == "pdf":
        return _extract_text_from_pdf_bytes(data)

    if file_type == "docx":
        return _extract_text_from_docx_bytes(data)

    if file_type == "doc":
        return _extract_text_from_legacy_doc_bytes(data)

    if file_type == "pptx":
        try:
            from pptx import Presentation

            prs = Presentation(BytesIO(data))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    text = getattr(shape, "text", None)
                    if text and str(text).strip():
                        parts.append(str(text).strip())
            return "\n".join(parts).strip()
        except Exception as e:
            raise RuntimeError(f"PPTX 텍스트 추출 실패: {e}") from e

    if file_type in ("xlsx", "xls"):
        try:
            import openpyxl

            wb = openpyxl.load_workbook(BytesIO(data), read_only=True, data_only=True)
            rows: list[str] = []
            for sheet in wb.worksheets:
                rows.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    vals = ["" if v is None else str(v) for v in row]
                    if any(v.strip() for v in vals):
                        rows.append("\t".join(vals))
            return "\n".join(rows).strip()
        except Exception as e:
            raise RuntimeError(f"Excel 텍스트 추출 실패: {e}") from e

    return ""


def _extract_text_from_local_document(file_path: str, file_type: str) -> str:
    """Best-effort text extraction for Load-files under /mnt/workspace."""
    with open(file_path, "rb") as f:
        data = f.read()
    return _extract_text_from_document_bytes(data, file_type)


def get_summary_of_uploaded_file(file_ref: str, prompt: str = "") -> str:
    """Analyze an uploaded file (by URL, workspace path, or name) and return a text summary.

    Images are loaded from S3 under images/{user_id}/ and summarized with vision.
    Load-files paths under /mnt/workspace/{user}/upload/ are read from the
    session-storage mount, with an S3 API fallback when the mount lags.
    """
    file_name = _file_name_from_ref(file_ref)
    if not file_name:
        return "파일 이름을 확인할 수 없습니다."

    file_type = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
    logger.info(f"get_summary_of_uploaded_file: file_name={file_name}, file_type={file_type}")

    # Load-files: agentcore-sessions → /mnt/workspace/{user}/upload/{name}
    local_path = (file_ref or "").strip()
    if local_path.startswith("/mnt/workspace/"):
        data, source = _load_workspace_file_bytes(local_path)
        if data is None:
            return (
                f"파일 경로: {local_path}\n"
                "파일을 마운트/S3에서 읽지 못했습니다. "
                f"read_file 또는 bash로 `{local_path}` 를 직접 읽어 분석하세요."
            )

        if file_type in ("png", "jpeg", "jpg", "webp", "gif"):
            image_summary_prompt = (
                "사용자의 요청을 참조하여 이미지의 내용을 분석한 후에 markdown 포맷으로 자세히 설명해주세요. "
                f"사용자 요청: <user_request>{prompt}</user_request>"
            )
            return summarize_image(data, image_summary_prompt)

        try:
            extracted = _extract_text_from_document_bytes(data, file_type)
        except Exception as e:
            logger.error(f"Failed to extract text from {local_path} ({source}): {e}")
            return (
                f"파일 경로: {local_path}\n"
                f"텍스트 추출에 실패했습니다 ({e}). "
                f"read_file / bash로 `{local_path}` 를 직접 읽어 분석하세요."
            )
        if not extracted:
            return (
                f"파일 경로: {local_path}\n"
                f"파일 형식(.{file_type})에서 텍스트를 추출하지 못했습니다. "
                f"read_file / bash로 `{local_path}` 를 직접 읽어 분석하세요."
            )
        max_chars = 100000
        if len(extracted) > max_chars:
            extracted = extracted[:max_chars] + "\n\n…(이하 생략)"
        return f"파일 경로: {local_path}\n\n{extracted}"

    if file_type in ("png", "jpeg", "jpg", "webp", "gif"):
        s3_client = boto3.client(
            service_name="s3",
            region_name=bedrock_region,
        )
        s3_key = _s3_key_from_file_ref(file_ref, default_prefix=s3_image_prefix)
        logger.info(f"loading image from s3://{s3_bucket}/{s3_key}")
        image_obj = s3_client.get_object(Bucket=s3_bucket, Key=s3_key)
        image_content = image_obj["Body"].read()

        image_summary_prompt = f"사용자의 요청을 참조하여 이미지의 내용을 분석한 후에 markdown 포맷으로 자세히 설명해주세요. 사용자 요청: <user_request>{prompt}</user_request>"
        logger.info(f"image_summary_prompt: {image_summary_prompt}")

        return summarize_image(image_content, image_summary_prompt)

    return f"지원하지 않는 파일 형식입니다: {file_type or '(unknown)'}"

####################### Bedrock Agent #######################
# RAG using Lambda
############################################################# 
def get_rag_prompt(text):
    # print("###### get_rag_prompt ######")
    llm = get_chat()
    # print('model_type: ', model_type)
    
    if model_type == "nova":
        if isKorean(text)==True:
            system = (
                "당신의 이름은 서연이고, 질문에 대해 친절하게 답변하는 사려깊은 인공지능 도우미입니다."
                "다음의 Reference texts을 이용하여 user의 질문에 답변합니다."
                "모르는 질문을 받으면 솔직히 모른다고 말합니다."
                "답변의 이유를 풀어서 명확하게 설명합니다."
            )
        else: 
            system = (
                "You will be acting as a thoughtful advisor."
                "Provide a concise answer to the question at the end using reference texts." 
                "If you don't know the answer, just say that you don't know, don't try to make up an answer."
                "You will only answer in text format, using markdown format is not allowed."
            )    
    
        human = (
            "Question: {question}"

            "Reference texts: "
            "{context}"
        ) 
        
    elif model_type == "claude":
        if isKorean(text)==True:
            system = (
                "당신의 이름은 서연이고, 질문에 대해 친절하게 답변하는 사려깊은 인공지능 도우미입니다."
                "다음의 <context> tag안의 참고자료를 이용하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다." 
                "모르는 질문을 받으면 솔직히 모른다고 말합니다."
                "답변의 이유를 풀어서 명확하게 설명합니다."
                "결과는 <result> tag를 붙여주세요."
            )
        else: 
            system = (
                "You will be acting as a thoughtful advisor."
                "Here is pieces of context, contained in <context> tags." 
                "If you don't know the answer, just say that you don't know, don't try to make up an answer."
                "You will only answer in text format, using markdown format is not allowed."
                "Put it in <result> tags."
            )    

        human = (
            "<question>"
            "{question}"
            "</question>"

            "<context>"
            "{context}"
            "</context>"
        )

    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    rag_chain = prompt | llm

    return rag_chain

bedrock_agent_runtime_client = boto3.client(
    "bedrock-agent-runtime",
    region_name=bedrock_region
)
knowledge_base_id = config.get('knowledge_base_id')
number_of_results = 4


def s3_uri_to_console_url(uri: str, region: str) -> str:
    """Open the object in the AWS S3 console (when sharing_url is not configured)."""
    if not uri or not uri.startswith("s3://"):
        return ""
    rest = uri[5:]
    parts = rest.split("/", 1)
    bucket = parts[0]
    key = parts[1] if len(parts) > 1 else ""
    enc_key = parse.quote(key, safe="")
    return f"https://{region}.console.aws.amazon.com/s3/object/{bucket}?prefix={enc_key}"


def retrieve(query):
    global knowledge_base_id
    
    try:
        response = bedrock_agent_runtime_client.retrieve(
            retrievalQuery={"text": query},
            knowledgeBaseId=knowledge_base_id,
                retrievalConfiguration={
                    "vectorSearchConfiguration": {"numberOfResults": number_of_results},
                },
            )
    except ClientError as e:
        error_code = e.response.get("Error", {}).get("Code", "")
        
        # Update knowledge_base_id only when ResourceNotFoundException occurs
        if error_code == "ResourceNotFoundException":
            logger.warning(f"ResourceNotFoundException occurred: {e}")
            logger.info("Attempting to update knowledge_base_id...")
            
            bedrock_region_local = config.get('region', 'us-west-2')
            projectName_local = config.get('projectName')

            # Create bedrock-agent client with same credentials as bedrock-agent-runtime client
            bedrock_agent_client = boto3.client("bedrock-agent", region_name=bedrock_region_local)
            knowledge_base_list = bedrock_agent_client.list_knowledge_bases()
            
            updated = False
            for knowledge_base in knowledge_base_list.get("knowledgeBaseSummaries", []):
                if knowledge_base["name"] == projectName_local:
                    new_knowledge_base_id = knowledge_base["knowledgeBaseId"]
                    knowledge_base_id = new_knowledge_base_id

                    config['knowledge_base_id'] = new_knowledge_base_id
                    with open(config_path, "w", encoding="utf-8") as f:
                        json.dump(config, f, ensure_ascii=False, indent=4)
                    
                    logger.info(f"Updated knowledge_base_id to: {new_knowledge_base_id}")
                    updated = True
                    break
            
            if updated:
                # Retry after updating knowledge_base_id
                try:
                    response = bedrock_agent_runtime_client.retrieve(
                        retrievalQuery={"text": query},
                        knowledgeBaseId=knowledge_base_id,
                        retrievalConfiguration={
                            "vectorSearchConfiguration": {"numberOfResults": number_of_results},
                        },
                    )
                    logger.info("Retry successful after updating knowledge_base_id")
                except Exception as retry_error:
                    logger.error(f"Retry failed after updating knowledge_base_id: {retry_error}")
                    raise
            else:
                logger.error(f"Could not find knowledge base with name: {projectName_local}")
                raise
        else:
            # Re-raise other errors that are not ResourceNotFoundException
            logger.error(f"Error retrieving: {e}")
            raise
    except Exception as e:
        # Re-raise other exceptions that are not ClientError
        logger.error(f"Unexpected error retrieving: {e}")
        raise
    
    # logger.info(f"response: {response}")
    retrieval_results = response.get("retrievalResults", [])
    # logger.info(f"retrieval_results: {retrieval_results}")

    json_docs = []
    for result in retrieval_results:
        text = url = name = None
        if "content" in result:
            content = result["content"]
            if "text" in content:
                text = content["text"]

        if "location" in result:
            location = result["location"]
            if "s3Location" in location:
                uri = location["s3Location"]["uri"] if location["s3Location"]["uri"] is not None else ""

                name = uri.split("/")[-1]
                encoded_name = parse.quote(name)
                if path:
                    url = f"{path}/{doc_prefix}{encoded_name}"
                else:
                    url = s3_uri_to_console_url(uri, bedrock_region)
                
            elif "webLocation" in location:
                url = location["webLocation"]["url"] if location["webLocation"]["url"] is not None else ""
                name = "WEB"

        page = None
        raw_page = (result.get("metadata") or {}).get("x-amz-bedrock-kb-document-page-number")
        if raw_page is not None:
            try:
                # Bedrock KB uses 0-based page numbers; convert to 1-based for display
                page = int(raw_page) + 1
            except (TypeError, ValueError):
                page = raw_page

        reference = {
            "url": url,
            "title": name,
            "from": "RAG",
        }
        if page is not None:
            reference["page"] = page

        json_docs.append({
            "contents": text,
            "reference": reference,
        })
    logger.info(f"json_docs: {json_docs}")

    return json.dumps(json_docs, ensure_ascii=False)
 
def run_rag_with_knowledge_base(query, st):
    global reference_docs, contentList
    reference_docs = []
    contentList = []

    # retrieve
    if debug_mode == "Enable":
        st.info(f"RAG 검색을 수행합니다. 검색어: {query}")  

    json_docs = retrieve(query)    
    logger.info(f"json_docs: {json_docs}")

    relevant_docs = json.loads(json_docs)

    relevant_context = ""
    for doc in relevant_docs:
        relevant_context += f"{doc['contents']}\n\n"

    # change format to document
    st.info(f"{len(relevant_docs)}개의 관련된 문서를 얻었습니다.")

    rag_chain = get_rag_prompt(query)
                       
    msg = ""    
    try: 
        result = rag_chain.invoke(
            {
                "question": query,
                "context": relevant_context                
            }
        )
        logger.info(f"result: {result}")

        msg = result.content        
        if msg.find('<result>')!=-1:
            msg = msg[msg.find('<result>')+8:msg.find('</result>')]        
               
    except Exception:
        err_msg = traceback.format_exc()
        logger.info(f"error message: {err_msg}")                    
        raise Exception ("Not able to request to LLM")
    
    if relevant_docs:
        refs = []
        for doc in relevant_docs:
            ref = {
                "title": doc["reference"]["title"],
                "url": doc["reference"]["url"],
                "content": doc["contents"],
            }
            if doc["reference"].get("page") is not None:
                ref["page"] = doc["reference"]["page"]
            refs.append(ref)
        msg += _format_references_markdown(refs)
    
    return msg, reference_docs
   
def extract_thinking_tag(response, st):
    if response.find('<thinking>') != -1:
        status = response[response.find('<thinking>')+10:response.find('</thinking>')]
        logger.info(f"gent_thinking: {status}")
        
        if debug_mode=="Enable":
            st.info(status)

        if response.find('<thinking>') == 0:
            msg = response[response.find('</thinking>')+12:]
        else:
            msg = response[:response.find('<thinking>')]
        logger.info(f"msg: {msg}")
    else:
        msg = response

    return msg

tool_input_list = dict()

sharing_url = config["sharing_url"] if "sharing_url" in config else None


def _urls_from_file_saved_message(tool_content) -> list:
    """Legacy/plain tool text: 'File saved: path' -> absolute paths for download UI."""
    text = tool_content
    if isinstance(tool_content, dict):
        text = tool_content.get("output", "") or ""
    if not isinstance(text, str) or "File saved:" not in text:
        return []
    tail = text.split("File saved:", 1)[1].strip()
    if not tail:
        return []
    line = tail.splitlines()[0].strip()
    full = line if os.path.isabs(line) else os.path.join(langgraph_agent.WORKING_DIR, line)
    full = os.path.normpath(full)
    if os.path.isfile(full):
        return [full]
    return []


def _parse_execute_code_artifact_paths(tool_content: str) -> list:
    """Parse absolute paths from execute_code output after an [artifacts] block."""
    if not isinstance(tool_content, str) or "[artifacts]" not in tool_content:
        return []
    idx = tool_content.find("[artifacts]")
    rest = tool_content[idx + len("[artifacts]") :].strip()
    out = []
    for line in rest.splitlines():
        line = line.strip()
        if line:
            out.append(line)
    return out


def _format_artifact_links_markdown(artifact_urls: list) -> str:
    """Append artifact list for the reply. Local files: relative path only (no file:// links)."""
    from pathlib import Path

    if not artifact_urls:
        return ""
    lines = ["", "### 생성된 파일"]
    for url in artifact_urls:
        name = url.split("/")[-1].split("?")[0]
        if not name or name == url:
            name = Path(url).name
        if url.startswith(("http://", "https://")):
            lines.append(f"- [{name}]({url})")
        else:
            try:
                rel = os.path.relpath(url, langgraph_agent.WORKING_DIR)
                rel = rel.replace("\\", "/")
            except (OSError, ValueError):
                rel = name
            lines.append(f"- `{rel}`")
    return "\n".join(lines) + "\n"
s3_prefix = "docs"
capture_prefix = "captures"

def _sanitize_reference_text(text: str, max_len: int) -> str:
    """Collapse whitespace/newlines and strip markdown that breaks list links."""
    if not text:
        return ""
    cleaned = " ".join(str(text).replace("\r", "\n").split())
    cleaned = cleaned.replace("```", "`").replace("[", "\\[").replace("]", "\\]")
    if len(cleaned) > max_len:
        cleaned = cleaned[: max_len - 3].rstrip(" .") + "..."
    return cleaned


def _format_references_markdown(references: list) -> str:
    """Build a Reference section safe for markdown list rendering."""
    lines = ["\n\n### Reference"]
    for i, reference in enumerate(references, start=1):
        title = _sanitize_reference_text(reference.get("title") or "Untitled", 120)
        content = _sanitize_reference_text(reference.get("content") or "", 100)
        url = (reference.get("url") or "").strip()
        page = reference.get("page")
        page_suffix = f" , {page} page" if page is not None else ""
        if url:
            lines.append(
                f"{i}. [{title}]({url}){page_suffix} — {content}" if content
                else f"{i}. [{title}]({url}){page_suffix}"
            )
        else:
            lines.append(
                f"{i}. {title}{page_suffix} — {content}" if content
                else f"{i}. {title}{page_suffix}"
            )
    return "\n".join(lines) + "\n"




def _build_tool_reference(ref_item: dict) -> dict:
    """Build a display reference from a RAG doc item."""
    reference = ref_item.get("reference") or {}
    contents = ref_item.get("contents") or ""
    content_text = contents[:100] + "..." if len(contents) > 100 else contents
    result = {
        "url": reference.get("url"),
        "title": reference.get("title"),
        "content": content_text,
    }
    if reference.get("page") is not None:
        result["page"] = reference["page"]
    return result


def get_tool_info(tool_name, tool_content):
    tool_references = []    
    urls = []
    content = ""
    
    # OpenSearch
    if tool_name == "SearchIndexTool": 
        if ":" in tool_content:
            extracted_json_data = tool_content.split(":", 1)[1].strip()
            try:
                json_data = json.loads(extracted_json_data)
                # logger.info(f"extracted_json_data: {extracted_json_data[:200]}")
            except json.JSONDecodeError:
                logger.info("JSON parsing error")
                json_data = {}
        else:
            json_data = {}
        
        if "hits" in json_data:
            hits = json_data["hits"]["hits"]
            if hits:
                logger.info(f"hits[0]: {hits[0]}")

            for hit in hits:
                text = hit["_source"]["text"]
                metadata = hit["_source"]["metadata"]
                
                content += f"{text}\n\n"

                filename = metadata["name"].split("/")[-1]
                # logger.info(f"filename: {filename}")
                
                content_part = text.replace("\n", "")
                tool_references.append({
                    "url": metadata["url"], 
                    "title": filename,
                    "content": content_part[:100] + "..." if len(content_part) > 100 else content_part
                })
                
        logger.info(f"content: {content}")
        
    # aws document
    elif tool_name == "search_documentation":
        try:
            # tool_content가 리스트인 경우 처리 (예: [{'type': 'text', 'text': '...'}])
            if isinstance(tool_content, list):
                # 리스트의 첫 번째 항목에서 text 필드 추출
                if len(tool_content) > 0 and isinstance(tool_content[0], dict) and 'text' in tool_content[0]:
                    tool_content = tool_content[0]['text']
                else:
                    logger.info(f"Unexpected list format: {tool_content}")
                    return content, urls, tool_references
            
            # tool_content가 문자열인 경우 JSON 파싱
            if isinstance(tool_content, str):
                json_data = json.loads(tool_content)
            elif isinstance(tool_content, dict):
                json_data = tool_content
            else:
                logger.info(f"Unexpected tool_content type: {type(tool_content)}")
                return content, urls, tool_references
            
            # search_results 배열에서 결과 추출
            search_results = json_data.get('search_results', [])
            if not search_results:
                # search_results가 없으면 json_data 자체가 배열일 수 있음
                if isinstance(json_data, list):
                    search_results = json_data
                else:
                    logger.info(f"No search_results found in JSON data")
                    return content, urls, tool_references
            
            for item in search_results:
                logger.info(f"item: {item}")
                
                if isinstance(item, str):
                    try:
                        item = json.loads(item)
                    except json.JSONDecodeError:
                        logger.info(f"Failed to parse item as JSON: {item}")
                        continue
                
                if isinstance(item, dict) and 'url' in item and 'title' in item:
                    url = item['url']
                    title = item['title']
                    content_text = item.get('context', '')[:100] + "..." if len(item.get('context', '')) > 100 else item.get('context', '')
                    tool_references.append({
                        "url": url,
                        "title": title,
                        "content": content_text
                    })
                else:
                    logger.info(f"Invalid item format: {item}")
                    
        except json.JSONDecodeError as e:
            logger.info(f"JSON parsing error: {e}, tool_content: {tool_content}")
            pass
        except Exception as e:
            logger.info(f"Unexpected error in search_documentation: {e}, tool_content type: {type(tool_content)}")
            pass

        logger.info(f"content: {content}")
        logger.info(f"tool_references: {tool_references}")
            
    # aws-knowledge
    elif tool_name == "aws___read_documentation":
        logger.info(f"#### {tool_name} ####")
        if isinstance(tool_content, dict):
            json_data = tool_content
        elif isinstance(tool_content, list):
            json_data = tool_content
        else:
            json_data = json.loads(tool_content)
        
        logger.info(f"json_data: {json_data}")
        payload = json_data["response"]["payload"]
        if "content" in payload:
            payload_content = payload["content"]
            if "result" in payload_content:
                result = payload_content["result"]
                logger.info(f"result: {result}")
                if isinstance(result, str) and "AWS Documentation from" in result:
                    logger.info(f"Processing AWS Documentation format: {result}")
                    try:
                        # Extract URL from "AWS Documentation from https://..."
                        url_start = result.find("https://")
                        if url_start != -1:
                            # Find the colon after the URL (not inside the URL)
                            url_end = result.find(":", url_start)
                            if url_end != -1:
                                # Check if the colon is part of the URL or the separator
                                url_part = result[url_start:url_end]
                                # If the colon is immediately after the URL, use it as separator
                                if result[url_end:url_end+2] == ":\n":
                                    url = url_part
                                    content_start = url_end + 2  # Skip the colon and newline
                                else:
                                    # Try to find the actual URL end by looking for space or newline
                                    space_pos = result.find(" ", url_start)
                                    newline_pos = result.find("\n", url_start)
                                    if space_pos != -1 and newline_pos != -1:
                                        url_end = min(space_pos, newline_pos)
                                    elif space_pos != -1:
                                        url_end = space_pos
                                    elif newline_pos != -1:
                                        url_end = newline_pos
                                    else:
                                        url_end = len(result)
                                    
                                    url = result[url_start:url_end]
                                    content_start = url_end + 1
                                
                                # Remove trailing colon from URL if present
                                if url.endswith(":"):
                                    url = url[:-1]
                                
                                # Extract content after the URL
                                if content_start < len(result):
                                    content_text = result[content_start:].strip()
                                    # Truncate content for display
                                    display_content = content_text[:100] + "..." if len(content_text) > 100 else content_text
                                    display_content = display_content.replace("\n", "")
                                    
                                    tool_references.append({
                                        "url": url,
                                        "title": "AWS Documentation",
                                        "content": display_content
                                    })
                                    content += content_text + "\n\n"
                                    logger.info(f"Extracted URL: {url}")
                                    logger.info(f"Extracted content length: {len(content_text)}")
                    except Exception as e:
                        logger.error(f"Error parsing AWS Documentation format: {e}")
        logger.info(f"content: {content}")
        logger.info(f"tool_references: {tool_references}")

    else:        
        try:
            if isinstance(tool_content, dict):
                json_data = tool_content
            elif isinstance(tool_content, list):
                json_data = tool_content
            else:
                json_data = json.loads(tool_content)
            
            logger.info(f"json_data: {json_data}")
            if isinstance(json_data, dict) and "path" in json_data:  # path
                path = json_data["path"]
                if isinstance(path, list):
                    for url in path:
                        urls.append(url)
                else:
                    urls.append(path)
            elif isinstance(json_data, list):  # Parse JSON from text field when json_data is a list
                for item in json_data:
                    if isinstance(item, dict) and "text" in item:
                        try:
                            text_json = json.loads(item["text"])
                            if isinstance(text_json, dict) and "path" in text_json:
                                path = text_json["path"]
                                if isinstance(path, list):
                                    for url in path:
                                        urls.append(url)
                                else:
                                    urls.append(path)
                        except (json.JSONDecodeError, TypeError):
                            pass            


            if isinstance(json_data, dict):
                for item in json_data:
                    logger.info(f"item: {item}")
                    if "reference" in item and "contents" in item:
                        tool_references.append(_build_tool_reference(item))
            elif isinstance(json_data, list):
                logger.info(f"json_data is a list: {json_data}")
                for item in json_data:
                    if isinstance(item, dict) and "text" in item:
                        try:
                            # text 필드 안의 JSON 문자열 파싱
                            text_json = json.loads(item["text"])
                            if isinstance(text_json, list):
                                # 파싱된 JSON이 리스트인 경우
                                for ref_item in text_json:
                                    if isinstance(ref_item, dict) and "reference" in ref_item and "contents" in ref_item:
                                        tool_references.append(_build_tool_reference(ref_item))
                            elif isinstance(text_json, dict) and "reference" in text_json and "contents" in text_json:
                                # 파싱된 JSON이 딕셔너리인 경우
                                tool_references.append(_build_tool_reference(text_json))
                        except (json.JSONDecodeError, TypeError) as e:
                            logger.warning(f"Failed to parse text JSON: {e}")
                            pass
                    elif isinstance(item, dict) and "reference" in item and "contents" in item:
                        # 리스트 항목이 직접 reference를 가지고 있는 경우
                        tool_references.append(_build_tool_reference(item))
                
            logger.info(f"tool_references: {tool_references}")

        except json.JSONDecodeError:
            pass

    if tool_name == "execute_code":
        extra: list = []
        if isinstance(tool_content, str):
            try:
                data = json.loads(tool_content)
                if isinstance(data, dict) and isinstance(data.get("output"), str):
                    extra.extend(_parse_execute_code_artifact_paths(data["output"]))
            except json.JSONDecodeError:
                extra.extend(_parse_execute_code_artifact_paths(tool_content))
        for u in extra:
            if u and u not in urls:
                urls.append(u)

    if not urls:
        extras: list = []
        if isinstance(tool_content, str):
            try:
                data = json.loads(tool_content)
                if isinstance(data, dict) and isinstance(data.get("output"), str):
                    extras.extend(_urls_from_file_saved_message(data["output"]))
            except json.JSONDecodeError:
                extras.extend(_urls_from_file_saved_message(tool_content))
        else:
            extras.extend(_urls_from_file_saved_message(tool_content))
        for u in extras:
            if u and u not in urls:
                urls.append(u)

    return content, urls, tool_references


def append_tool_guidance_to_prompt(system_prompt: str, mcp_servers: list) -> str:
    """Append tool-specific usage guidance to the system prompt based on selected MCP servers.

    Add new guidance rules here as MCP servers are introduced.
    """
    if not mcp_servers:
        return system_prompt

    selected = {name.lower() for name in mcp_servers}
    extras: list[str] = []

    if "memory" in selected:
        extras.append(skill.MEMORY_RECALL_GUIDANCE)

    parallel_tools: list[str] = []
    if "wiki" in selected:
        parallel_tools.append("recall_wiki")
    if "tavily" in selected:
        parallel_tools.append("tavily_web_search")
    if "websearch" in selected:
        parallel_tools.append("websearch")
    if "knowledge base" in selected:
        parallel_tools.append("retrieve")
    if len(parallel_tools) >= 2:
        extras.append(
            f"검색이 필요한 경우에 {', '.join(parallel_tools)}을 이용해 병렬로 조회하세요."
        )

    if "aws documentation" in selected:
        extras.append(
            "aws와 관련된 내용이 있다면, search_documentation tool을 이용해 필요한 정보를 수집하세요."
        )

    if not extras:
        return system_prompt

    logger.info(f"extra prompt: {extras}")

    return system_prompt + "\n" + "\n".join(extras)


async def create_agent(
    mcp_servers: list,
    skill_list: list,
    runtime_session_id: str | None = None,
) -> tuple[str, list]:
    session_id = runtime_session_id or _checkpoint_session_id()
    set_checkpoint_session_id(session_id)
    thread_scope = _thread_scope(session_id)

    await ensure_checkpointer()

    # builtin tools
    tools = langgraph_agent.get_builtin_tools()
    logger.info(f"builtin_tools count: {len(tools)}")
        
    # mcp
    mcp_json = mcp_config.load_selected_config(mcp_servers)
    # logger.info(f"mcp_json: {mcp_json}")

    server_params = langgraph_agent.load_multiple_mcp_server_parameters(mcp_json)

    # Pass current user_id to per-user MCP servers via process env.
    # UI may send "knowledge base"; load_selected_config maps it to "kb-retriever".
    scoped_user_id = user_id if user_id else "default"
    for server_name in (
        "memory",
        "kb-retriever",
        "knowledge base",
        "imageGeneration",
        "image_generation",
    ):
        params = server_params.get(server_name)
        if params and params.get("transport") == "stdio":
            env = dict(params.get("env") or {})
            env["AGENTCORE_USER_ID"] = scoped_user_id
            params["env"] = env
            logger.info(f"{server_name} MCP AGENTCORE_USER_ID={scoped_user_id}")

    has_agentcore = any(
        cfg.get("auth_type") == "aws_sigv4"
        for cfg in (mcp_json.get("mcpServers") or {}).values()
    )

    if server_params:
        try:
            mcp_tools = []
            max_retries = 3 if has_agentcore else 1
            if has_agentcore:
                logger.info(
                    "Loading MCP tools from Bedrock AgentCore (cold start may take 1-2 minutes)..."
                )
            for server_name, params in server_params.items():
                loaded = False
                for attempt in range(1, max_retries + 1):
                    try:
                        async with MCPAdapter({"mcpServers": {server_name: params}}) as adapter:
                            logger.info(f"MCP client initialized for server: {server_name}")
                            _tools = await adapter.list_tools()
                        mcp_tools.extend(_tools)
                        loaded = True
                        break
                    except Exception as e:
                        if attempt >= max_retries:
                            logger.error(f"Failed to load MCP server '{server_name}': {e}")
                            break
                        wait_seconds = attempt * 15
                        logger.warning(
                            "MCP list_tools attempt %s/%s for %s failed: %s. Retrying in %ss...",
                            attempt,
                            max_retries,
                            server_name,
                            e,
                            wait_seconds,
                        )
                        await asyncio.sleep(wait_seconds)

            mcp_tools = wrap_tavily_mcp_tools(mcp_tools)

            for tool in mcp_tools:
                logger.info(f"mcp_tool: {tool.name}")
                if tool.name not in [t.name for t in tools]:
                    tools.append(tool)
                else:
                    logger.info(f"mcp_tool of {tool.name} already in tools")
        except Exception as e:
            logger.error(f"Error creating MCP client or getting tools: {e}")
            if getattr(e, "__cause__", None):
                logger.error(f"  cause: {e.__cause__}")
            logger.info(f"Falling back to builtin tools only (count: {len(tools)})")

    if "aws-tavily" in mcp_servers and not any(
        getattr(t, "name", "").startswith("tavily_") for t in tools
    ):
        logger.warning(
            "aws-tavily was selected but no tavily_* tools were loaded. "
            "Check agent_runtime_aws_tavily in us-east-1 and Runtime IAM permissions."
        )

    tools.extend(skill.get_skill_tools())

    skill_info = skill.get_skill_info(skill_list)
    logger.info(f"skill_info: {skill_info}")

    system_prompt = skill.build_skill_prompt(skill_info)
    if any(getattr(t, "name", "").startswith("tavily_") for t in tools):
        system_prompt += langgraph_agent.TAVILY_TOOL_PROMPT

    tool_list = [tool.name for tool in tools] if tools else []
    logger.info(f"tool_list: {tool_list}")
    system_prompt = append_tool_guidance_to_prompt(system_prompt, mcp_servers)

    if not tools:
        logger.warning("No tools available, using general conversation mode")
        return None, None
    
    thread_id = thread_scope

    active_checkpointer = await ensure_checkpointer()
    checkpoint_db = get_checkpoint_db_path()
    logger.info(
        "checkpoint bind: session_id=%s thread_id=%s db=%s saver=%s",
        session_id,
        thread_id,
        checkpoint_db,
        type(active_checkpointer).__name__,
    )

    app = langgraph_agent.buildChatAgentWithHistory(tools, checkpointer=active_checkpointer)
    agent_config = {
        "recursion_limit": 100,
        "configurable": {
            "thread_id": thread_id,
            "tools": tools,
            "system_prompt": system_prompt,
        },
        "max_turns": langgraph_agent.MAX_CONTEXT_TURNS,
    }
    
    return app, agent_config

app = None
agent_config = None
active_mcp_servers = []
active_skills = []
current_id = None
_active_agent_session = None


async def get_or_create_agent(
    mcp_servers: list,
    skill_list: list,
    runtime_session_id: str | None = None,
):
    """Reuse the previously built LangGraph app/tools when nothing changed.

    ``create_agent()`` always rebuilds every MCP client from scratch (which
    respawns stdio subprocesses such as memory/wiki/kb-retriever). That is
    correct the first time a microVM handles a session, but on every
    subsequent invoke within the *same* microVM + session, rebuilding is
    pure waste if mcp_servers/skill_list/user_id/session are unchanged.

    This wrapper compares the current call's parameters against the last
    successful build (kept in module-level globals, i.e. scoped to this
    microVM's process) and short-circuits to the cached (app, agent_config)
    when nothing relevant changed.
    """
    global app, agent_config, active_mcp_servers, active_skills, current_id, _active_agent_session

    session_id = runtime_session_id or _checkpoint_session_id()

    if (
        app is not None
        and active_mcp_servers == mcp_servers
        and active_skills == skill_list
        and current_id == user_id
        and _active_agent_session == session_id
    ):
        logger.info(
            f"Reusing cached agent/MCP tools (session={session_id}, "
            f"mcp_servers={mcp_servers}, skills={skill_list})"
        )
        return app, agent_config

    logger.info(
        f"Building new agent/MCP tools (session changed: {_active_agent_session} -> {session_id}, "
        f"or mcp_servers/skills/user_id changed)"
    )
    active_mcp_servers = mcp_servers
    active_skills = skill_list
    current_id = user_id
    _active_agent_session = session_id

    app, agent_config = await create_agent(mcp_servers, skill_list, runtime_session_id)
    return app, agent_config


async def run_langgraph_agent(query: str, mcp_servers: list, skill_list: list):
    global app, agent_config, active_mcp_servers, active_skills, current_id

    artifacts = []
    references = []

    if app is None or active_mcp_servers != mcp_servers or active_skills != skill_list or current_id != user_id:
        active_mcp_servers = mcp_servers
        active_skills = skill_list
        current_id = user_id

        app, agent_config = await create_agent(mcp_servers, skill_list)
    
    if app is None:
        logger.error("Failed to create agent - app is None")
        return "에이전트를 생성할 수 없습니다. MCP 서버 설정 또는 도구 구성을 확인해주세요.", []
    
    inputs = {
        "messages": [HumanMessage(content=query)]
    }
            
    result = ""
    tool_used = False  # Track if tool was used
    tool_name = toolUseId = ""
    async for stream in app.astream(inputs, agent_config, stream_mode="messages"):
        if isinstance(stream[0], AIMessageChunk):
            message = stream[0]    
            input = {}        
            if isinstance(message.content, list):
                for content_item in message.content:
                    if isinstance(content_item, dict):
                        if content_item.get('type') == 'text':
                            text_content = content_item.get('text', '')
                            # logger.info(f"text_content: {text_content}")
                            
                            # If tool was used, start fresh result
                            if tool_used:
                                result = text_content
                                tool_used = False
                            else:
                                result += text_content
                                
                            # logger.info(f"result: {result}")

                        elif content_item.get('type') == 'tool_use':
                            # logger.info(f"content_item: {content_item}")      
                            if 'id' in content_item and 'name' in content_item:
                                toolUseId = content_item.get('id', '')
                                tool_name = content_item.get('name', '')
                                logger.info(f"tool_name: {tool_name}, toolUseId: {toolUseId}")

                            if 'partial_json' in content_item:
                                partial_json = content_item.get('partial_json', '')
                                
                                if toolUseId not in tool_input_list:
                                    tool_input_list[toolUseId] = ""                                
                                tool_input_list[toolUseId] += partial_json
                                input = tool_input_list[toolUseId]
                        
        elif isinstance(stream[0], ToolMessage):
            message = stream[0]
            logger.info(f"ToolMessage: {message.name}, {message.content}")
            tool_name = message.name
            toolResult = message.content
            toolUseId = message.tool_call_id
            logger.info(f"toolResult: {toolResult}, toolUseId: {toolUseId}")
            tool_used = True
            
            content, urls, refs = get_tool_info(tool_name, toolResult)
            if refs:
                for r in refs:
                    references.append(r)
                logger.info(f"refs: {refs}")
            if urls:
                for url in urls:
                    artifacts.append(url)
                logger.info(f"urls: {urls}")

            if content:
                logger.info(f"content: {content}")        
    
    if not result:
        result = "답변을 찾지 못하였습니다."        
    logger.info(f"result: {result}")

    if references:
        result += _format_references_markdown(references)

    if artifacts:
        result += _format_artifact_links_markdown(artifacts)

    return result, artifacts


#########################################################
# AgentCore Memory
#########################################################
import agentcore_memory


def initiate_memory():
    """Load or create AgentCore Memory session variables for the current user."""
    global memory_id, actor_id, session_id

    effective_user_id = user_id if user_id and str(user_id).strip() else "default"
    logger.info(f"initiate_memory for user_id: {effective_user_id}")

    memory_id, actor_id, session_id, namespace = agentcore_memory.load_memory_variables(
        effective_user_id
    )
    # actor_id is the sanitized/aliased memory identity (not raw email)
    if not namespace:
        namespace = f"/users/{actor_id}/preferences"
    logger.info(
        f"memory_id: {memory_id}, actor_id: {actor_id}, "
        f"session_id: {session_id}, namespace: {namespace}"
    )

    if memory_id is None:
        memory_id = agentcore_memory.retrieve_memory_id()
        if memory_id is None:
            logger.info("Memory will be created...")
            memory_id = agentcore_memory.create_memory()
            logger.info(f"Memory was created... {memory_id}")

    # Shared strategies per memory_id (UserPreference / Summary / Semantic).
    # Namespace templates use {actorId}/{sessionId} — not one strategy per user.
    agentcore_memory.create_strategy_if_not_exists(memory_id)


def save_to_memory(query, result):
    """Save conversation to AgentCore Memory when memory_enabled is True."""
    global memory_id, actor_id, session_id

    if not memory_enabled:
        return

    try:
        expected_actor = agentcore_memory.resolve_memory_actor_id(
            user_id if user_id and str(user_id).strip() else "default"
        )
        if memory_id is None or actor_id != expected_actor:
            initiate_memory()

        agentcore_memory.save_conversation_to_memory(
            memory_id, actor_id, session_id, query, result
        )
        logger.info(f"Saved conversation to AgentCore Memory for actor_id={actor_id}")
    except Exception as e:
        logger.error(f"Failed to save conversation to AgentCore Memory: {e}")
